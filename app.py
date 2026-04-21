from flask import Flask, render_template, request, jsonify, session, send_file
import ollama
import PyPDF2
import psycopg2
from pptx import Presentation
from dotenv import load_dotenv
import os
from diffusers import StableDiffusionPipeline
import torch
import time
import tempfile



# ---------------- LOAD ENV ----------------
load_dotenv()
app = Flask(__name__)
app.secret_key = os.getenv("SECRET_KEY")
# ---------------- DATABASE ----------------
def get_conn():
    return psycopg2.connect(
        host=os.getenv("DB_HOST"),
        database=os.getenv("DB_NAME"),
        user=os.getenv("DB_USER"),
        password=os.getenv("DB_PASSWORD")
    )

# ---------------- IMAGE MODEL ----------------
device = "cuda" if torch.cuda.is_available() else "cpu"

pipe = StableDiffusionPipeline.from_pretrained(
    "runwayml/stable-diffusion-v1-5",
    torch_dtype=torch.float16 if device == "cuda" else torch.float32
).to(device)

# ---------------- HELPERS ----------------
def generate_title(text):
    return " ".join(text.split()[:5]) if text else "New Chat"

def update_title(chat_id, text):
    try:
        conn = get_conn()
        cur = conn.cursor()
        cur.execute("UPDATE chats SET title=%s WHERE id=%s",
                    (generate_title(text), chat_id))
        conn.commit()
        cur.close()
        conn.close()
    except:
        pass

def save_message(chat_id, user_msg, ai_msg):
    try:
        conn = get_conn()
        cur = conn.cursor()
        cur.execute(
            "INSERT INTO messages (chat_id, user_msg, ai_msg) VALUES (%s, %s, %s)",
            (chat_id, user_msg, ai_msg)
        )
        conn.commit()
        cur.close()
        conn.close()
    except:
        pass

def get_context(chat_id, limit=5):
    conn = get_conn()
    cur = conn.cursor()

    cur.execute(
        "SELECT user_msg, ai_msg FROM messages WHERE chat_id=%s ORDER BY id DESC LIMIT %s",
        (chat_id, limit)
    )

    rows = cur.fetchall()[::-1]
    cur.close()
    conn.close()

    messages = [{"role": "system", "content": "You are a helpful AI assistant."}]

    for u, a in rows:
        if not a.startswith("[IMAGE]"):
            messages.append({"role": "user", "content": u})
            messages.append({"role": "assistant", "content": a})

    return messages

# ---------------- ROUTES ----------------
@app.route("/")
def home():
    return render_template("index.html")

# ---------------- CHAT ----------------
@app.route("/chat", methods=["POST"])
def chat():
    try:
        user_input = request.json.get("message", "")
        chat_id = session.get("chat_id")

        conn = get_conn()
        cur = conn.cursor()

        if not chat_id:
            cur.execute("INSERT INTO chats (title) VALUES (%s) RETURNING id", ("New Chat",))
            chat_id = cur.fetchone()[0]
            session["chat_id"] = chat_id
            conn.commit()

        # -------- IMAGE --------
        if any(k in user_input.lower() for k in ["generate image", "draw", "image of"]):
            prompt = user_input.lower()
            for k in ["generate image", "draw", "image of"]:
                prompt = prompt.replace(k, "")

            image = pipe(prompt.strip()).images[0]

            os.makedirs("static/images", exist_ok=True)
            filename = f"{chat_id}_{int(time.time())}.png"
            path = f"static/images/{filename}"
            image.save(path)

            save_message(chat_id, user_input, f"[IMAGE]{path}")
            update_title(chat_id, user_input)

            cur.close()
            conn.close()

            return jsonify({
                "response": "Image generated 🎨",
                "image_url": "/" + path
            })

        # -------- CONTEXT --------
        # -------- FILE CONTEXT --------
        conn2 = get_conn()
        cur2 = conn2.cursor()

        cur2.execute("SELECT file_context FROM chats WHERE id=%s", (chat_id,))
        row = cur2.fetchone()

        file_context = row[0] if row and row[0] else ""

        cur2.close()
        conn2.close()

        messages = get_context(chat_id)

# Inject file context ONLY if exists
        if file_context:
                messages.append({
                    "role": "system",
                    "content": f"Use this document to answer:\n{file_context[:3000]}"
                })

        messages.append({"role": "user", "content": user_input})

        response = ollama.chat(
            model="llama3",
            messages=messages
        )

        ai_reply = response["message"]["content"]

        save_message(chat_id, user_input, ai_reply)
        update_title(chat_id, user_input)

        cur.close()
        conn.close()

        return jsonify({"response": ai_reply})

    except Exception as e:
        print("CHAT ERROR:", e)
        return jsonify({"response": "Error ❌"})

# ---------------- DOWNLOAD ----------------
@app.route("/download/<int:chat_id>")
def download(chat_id):
    conn = get_conn()
    cur = conn.cursor()

    cur.execute("SELECT user_msg, ai_msg FROM messages WHERE chat_id=%s", (chat_id,))
    rows = cur.fetchall()

    cur.close()
    conn.close()

    temp = tempfile.NamedTemporaryFile(delete=False, suffix=".txt")
    with open(temp.name, "w", encoding="utf-8") as f:
        for u, a in rows:
            f.write(f"User: {u}\nAI: {a}\n\n")

    return send_file(temp.name, as_attachment=True, download_name="chat.txt")

# ---------------- RENAME ----------------
@app.route("/rename_chat", methods=["POST"])
def rename_chat():
    data = request.json

    conn = get_conn()
    cur = conn.cursor()

    cur.execute("UPDATE chats SET title=%s WHERE id=%s",
                (data["title"], data["chat_id"]))
    conn.commit()

    cur.close()
    conn.close()

    return jsonify({"status": "ok"})

# ---------------- CLEAR ----------------
@app.route("/clear_chat/<int:chat_id>", methods=["DELETE"])
def clear_chat(chat_id):
    conn = get_conn()
    cur = conn.cursor()

    cur.execute("DELETE FROM messages WHERE chat_id=%s", (chat_id,))
    conn.commit()

    cur.close()
    conn.close()

    return jsonify({"status": "cleared"})

# ---------------- FILE UPLOAD ----------------
@app.route("/upload", methods=["POST"])
def upload():
    try:
        file = request.files["file"]
        filename = file.filename.lower()
        text = ""

        if filename.endswith(".pdf"):
            reader = PyPDF2.PdfReader(file)
            for p in reader.pages:
                text += p.extract_text() or ""

        elif filename.endswith(".txt"):
            text = file.read().decode()

        elif filename.endswith(".pptx"):
            prs = Presentation(file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"

        if not text.strip():
            return jsonify({"response": "No text found ❌"})

        chat_id = session.get("chat_id")

        if chat_id:
            conn = get_conn()
            cur = conn.cursor()

            cur.execute("UPDATE chats SET file_context=%s WHERE id=%s",
                        (text, chat_id))
            conn.commit()

            cur.close()
            conn.close()

        return jsonify({"response": "File uploaded ✅"})

    except Exception as e:
        print("UPLOAD ERROR:", e)
        return jsonify({"response": "Upload failed ❌"})

# ---------------- GET CHATS ----------------
@app.route("/get_chats")
def get_chats():
    conn = get_conn()
    cur = conn.cursor()

    cur.execute("SELECT id, title FROM chats ORDER BY id DESC")
    rows = cur.fetchall()

    cur.close()
    conn.close()

    return jsonify([{"id": r[0], "title": r[1]} for r in rows])

# ---------------- GET MESSAGES ----------------
@app.route("/get_messages/<int:chat_id>")
def get_messages(chat_id):
    session["chat_id"] = chat_id

    conn = get_conn()
    cur = conn.cursor()

    cur.execute("SELECT user_msg, ai_msg FROM messages WHERE chat_id=%s ORDER BY id", (chat_id,))
    rows = cur.fetchall()

    cur.close()
    conn.close()

    result = []

    for u, a in rows:
        if a.startswith("[IMAGE]"):
            result.append({"user": u, "image": "/" + a.replace("[IMAGE]", "")})
        else:
            result.append({"user": u, "bot": a})

    return jsonify(result)

# ---------------- NEW CHAT ----------------
@app.route("/new_chat", methods=["POST"])
def new_chat():
    conn = get_conn()
    cur = conn.cursor()

    cur.execute("INSERT INTO chats (title) VALUES (%s) RETURNING id", ("New Chat",))
    chat_id = cur.fetchone()[0]
    conn.commit()

    session["chat_id"] = chat_id

    cur.close()
    conn.close()

    return jsonify({"chat_id": chat_id})

# ---------------- DELETE ----------------
@app.route("/delete_chat/<int:chat_id>", methods=["DELETE"])
def delete_chat(chat_id):
    conn = get_conn()
    cur = conn.cursor()

    cur.execute("DELETE FROM messages WHERE chat_id=%s", (chat_id,))
    cur.execute("DELETE FROM chats WHERE id=%s", (chat_id,))
    conn.commit()

    cur.close()
    conn.close()

    return jsonify({"status": "deleted"})

# ---------------- RUN ----------------
if __name__ == "__main__":
    app.run(debug=False)